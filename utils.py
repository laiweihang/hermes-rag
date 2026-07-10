"""文档加载与解析工具集。

定位：把任意业务文件（PDF/Word/Excel/PPT/HTML/CSV/JSON/图片）统一转
换成 ``langchain_core.documents.Document`` 列表，供下游切片+嵌入使用。

入口：``load_document(file_path)``。其余 ``_xxx`` 私有函数按格式分派。

设计要点：
- ``UPLOAD_DIR`` 仅用于占位 import 兼容；本模块不直接读 uploads 目录，
  调用方传完整路径。
- 重型依赖（openpyxl / python-pptx / bs4 / ocr_engine）按需 lazy import，
  避免冷启动加载所有库。
- 所有解析失败统一返回空列表 ``[]``，由上层决定是否提示用户。
"""

import json
import logging
import os
import shutil
import subprocess
from langchain_core.documents import Document
from langchain_community.document_loaders import (
    PyPDFLoader,
    Docx2txtLoader,
    TextLoader,
    CSVLoader,
)
from config import UPLOAD_DIR, IMAGE_EXTENSIONS, OCR_MIN_TEXT_THRESHOLD

logger = logging.getLogger(__name__)


def _is_scanned_pdf(docs) -> bool:
    """判断 PyPDFLoader 提取出的文档是否为扫描件（大部分页面文本极少）。

    判定逻辑：超过 50% 的页面字符数低于 ``OCR_MIN_TEXT_THRESHOLD`` 即视
    为扫描件。允许少量正常文本页是因为很多扫描件首页/末页是封面 + 目录，
    可能含少量原生文本。
    """
    if not docs:
        # PyPDF 完全提取不到内容 —— 当作扫描件让上层走 OCR。
        return True
    scanned_pages = sum(
        1 for d in docs if len(d.page_content.strip()) < OCR_MIN_TEXT_THRESHOLD
    )
    return scanned_pages > len(docs) * 0.5


def _convert_doc_to_docx(file_path: str) -> str | None:
    """尝试用本机 LibreOffice 把 .doc 转成 .docx，返回新文件绝对路径；
    找不到 LibreOffice 时返回 None。

    选择 LibreOffice 而非纯 Python 库（如 antiword / textract）是因为：
    .doc 是二进制复合格式，纯 Python 解析对带表格、嵌入对象的文档兼容
    性差；LibreOffice 是事实上最完善的开源转换器。

    Windows 用户可在 https://www.libreoffice.org/ 下载安装；
    安装后 PATH 里会出现 ``soffice`` 命令。
    """
    # 跨平台命令名探测：Linux 多叫 libreoffice，macOS/Windows 叫 soffice。
    soffice = shutil.which("soffice") or shutil.which("libreoffice")
    if not soffice:
        return None
    out_dir = os.path.dirname(os.path.abspath(file_path))
    try:
        subprocess.run(
            [
                soffice,
                "--headless",          # 无 GUI 模式
                "--convert-to",
                "docx",
                "--outdir",
                out_dir,
                file_path,
            ],
            check=True,
            capture_output=True,
            timeout=120,               # 大文档最多给 2 分钟
        )
    except Exception as e:
        logger.error(f"LibreOffice 转换 .doc 失败: {e}")
        return None
    new_path = os.path.splitext(file_path)[0] + ".docx"
    # 进程退出码为 0 不等于一定生成了文件，再做一次存在性确认。
    return new_path if os.path.exists(new_path) else None


def _load_html(file_path: str):
    """读取 HTML 文件，BeautifulSoup 提取纯文本。"""
    from bs4 import BeautifulSoup

    # errors="ignore" 容忍编码异常 —— 网页常见声明 UTF-8 实际混入 GBK 字节。
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        soup = BeautifulSoup(f.read(), "lxml")
    # 去掉脚本/样式/noscript：纯文本对 RAG 有用，代码不利于检索。
    for tag in soup(["script", "style", "noscript"]):
        tag.decompose()
    text = soup.get_text("\n", strip=True)
    return [Document(page_content=text, metadata={"source": file_path})]


def _load_xlsx(file_path: str):
    """读 xlsx：每个 sheet 拼成一段文本（CSV 风格），返回一个 Document。

    用 read_only=True 避免大表整体加载到内存；data_only=True 让公式单元
    返回上次计算结果而不是公式字符串。
    """
    from openpyxl import load_workbook

    wb = load_workbook(file_path, data_only=True, read_only=True)
    parts = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            # 跳过整行为空的行 —— Excel 末尾常有大段空行虚增 chunk 数。
            if any(c.strip() for c in cells):
                rows.append("\t".join(cells))
        if rows:
            # Sheet 名加在前面，便于 LLM 引用时区分多个工作表。
            parts.append(f"# Sheet: {sheet.title}\n" + "\n".join(rows))
    text = "\n\n".join(parts) if parts else ""
    return [Document(page_content=text, metadata={"source": file_path})]


def _load_pptx(file_path: str):
    """读 pptx：每张幻灯片合成一段文本。"""
    from pptx import Presentation

    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            # 表格、图片等 shape 没有 text 属性，hasattr 判断避免 AttributeError。
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
        if texts:
            slides.append(f"--- Slide {i} ---\n" + "\n".join(texts))
    text = "\n\n".join(slides) if slides else ""
    return [Document(page_content=text, metadata={"source": file_path})]


def _load_json(file_path: str):
    """读 JSON：以缩进字符串形式作为正文。

    缩进序列化是为了让向量切片器能在结构边界（key/value 行）处优先切断，
    比 minified JSON 检索效果好。解析失败时退回原始文本，兼容 JSONL 等
    非标准 JSON 内容。
    """
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        try:
            data = json.load(f)
            text = json.dumps(data, ensure_ascii=False, indent=2)
        except Exception:
            f.seek(0)
            text = f.read()
    return [Document(page_content=text, metadata={"source": file_path})]


def load_document(file_path: str):
    """根据文件后缀加载文档；扫描件 PDF 与图片走 OCR；.doc 自动转 .docx。

    返回值：``list[Document]``。失败时返回空列表（不抛异常），由上层
    决定是否提示用户重试。

    分派策略：每种格式对应一个私有 loader；ext 列表见函数体。
    新增格式时：在 ``config.ALLOWED_UPLOAD_EXTENSIONS`` 加扩展名 + 在
    本函数加 ``if ext == ...`` 分支。
    """
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".pdf":
            docs = PyPDFLoader(file_path).load()
            # 文本提取后立即判断是不是扫描件，决定是否回退 OCR。
            if _is_scanned_pdf(docs):
                logger.info(
                    f"[OCR] 检测到扫描件 PDF，切换到 OCR 模式: {os.path.basename(file_path)}"
                )
                # lazy import：未启用 OCR 时避免拉起 GLM-OCR 依赖。
                from ocr_engine import ocr_pdf

                docs = ocr_pdf(file_path)
            return docs

        if ext in IMAGE_EXTENSIONS:
            from ocr_engine import ocr_image_file

            return ocr_image_file(file_path)

        if ext == ".docx":
            return Docx2txtLoader(file_path).load()

        if ext == ".doc":
            # .doc 必须先转 .docx；未装 LibreOffice 时给出清晰的错误指引。
            converted = _convert_doc_to_docx(file_path)
            if converted is None:
                raise RuntimeError(
                    "无法解析 .doc：未检测到 LibreOffice (soffice)。"
                    "请将文件在 Word 中另存为 .docx 后重新上传，"
                    "或安装 LibreOffice (https://www.libreoffice.org/) 让后端自动转换。"
                )
            logger.info(f"[DOC] 已通过 LibreOffice 转换: {os.path.basename(file_path)} -> {os.path.basename(converted)}")
            return Docx2txtLoader(converted).load()

        if ext in (".txt", ".md", ".markdown"):
            return TextLoader(file_path, encoding="utf-8").load()

        if ext in (".html", ".htm"):
            return _load_html(file_path)

        if ext == ".csv":
            return CSVLoader(file_path, encoding="utf-8").load()

        if ext == ".json":
            return _load_json(file_path)

        if ext == ".xlsx":
            return _load_xlsx(file_path)

        if ext == ".pptx":
            return _load_pptx(file_path)

        raise ValueError(f"不支持的文件格式: {ext}")
    except Exception as e:
        # 统一吞异常返回空列表 —— 上层 ingest / api 用 ``not docs`` 判断
        # 是否需要提示"未提取到内容"。
        logger.error(f"加载文件失败 {file_path}: {e}")
        return []
